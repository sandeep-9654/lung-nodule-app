#!/usr/bin/env python3
"""Convert presentation.html slides to a .pptx file."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ---------- constants ----------
W, H = Inches(13.333), Inches(7.5)  # widescreen 16:9
BG    = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0x4F, 0x46, 0xE5)
ACCENT2 = RGBColor(0x63, 0x66, 0xF1)
TEAL   = RGBColor(0x0D, 0x94, 0x88)
TEXT   = RGBColor(0x0F, 0x17, 0x2A)
TEXT2  = RGBColor(0x47, 0x55, 0x69)
TEXT3  = RGBColor(0x94, 0xA3, 0xB8)
AMBER  = RGBColor(0xD9, 0x77, 0x06)
ROSE   = RGBColor(0xE1, 0x1D, 0x48)
SKY    = RGBColor(0x02, 0x84, 0xC7)
EMERALD = RGBColor(0x05, 0x96, 0x69)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)

FONT   = "Inter"
DIR    = os.path.dirname(os.path.abspath(__file__))

prs = Presentation()
prs.slide_width = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]  # blank layout


def add_bg(slide, color=BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def tb(slide, left, top, width, height, text, size=18, bold=False, color=TEXT, align=PP_ALIGN.LEFT, font=FONT):
    """Add a textbox and return (textbox, textframe)."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font
    p.alignment = align
    return txBox, tf


def add_para(tf, text, size=18, bold=False, color=TEXT2, align=PP_ALIGN.LEFT, space_before=Pt(6)):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = FONT
    p.alignment = align
    p.space_before = space_before
    return p


def add_bullet(tf, text, size=16, color=TEXT2, level=0):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.name = FONT
    p.level = level
    p.space_before = Pt(4)
    return p


def add_accent_header(slide, num, title, top=Inches(0.6)):
    _, tf = tb(slide, Inches(0.8), top, Inches(11), Inches(0.7),
               f"{num:02d}.", size=28, bold=True, color=ACCENT, align=PP_ALIGN.LEFT)
    run = tf.paragraphs[0].add_run()
    run.text = f"  {title}"
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = TEXT
    run.font.name = FONT


def add_card(slide, left, top, width, height, title, body, icon_color=ACCENT, title_size=16, body_size=14):
    """Draw a card-like rounded rectangle with title and body text."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xF8, 0xFA, 0xFC)
    shape.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(14)
    tf.margin_right = Pt(14)
    tf.margin_top = Pt(10)
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(title_size)
    p.font.bold = True
    p.font.color.rgb = TEXT
    p.font.name = FONT
    if body:
        bp = tf.add_paragraph()
        bp.text = body
        bp.font.size = Pt(body_size)
        bp.font.color.rgb = TEXT2
        bp.font.name = FONT
        bp.space_before = Pt(6)
    return shape


# ============================================================
# SLIDE 1 — INTRODUCTION + TEAM
# ============================================================
s = prs.slides.add_slide(BLANK); add_bg(s)
tb(s, Inches(1), Inches(0.4), Inches(11), Inches(0.4),
   "B.Tech Final Year Project — CSE (Data Science)", size=14, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
tb(s, Inches(1), Inches(1.0), Inches(11), Inches(1.4),
   "Volumetric CT Scan Analysis for Pulmonary\nNodule Detection and Malignancy Classification Using\n3D CNN Architectures",
   size=30, bold=True, color=TEXT, align=PP_ALIGN.CENTER)
tb(s, Inches(2), Inches(2.6), Inches(9), Inches(0.7),
   "A smart system that automatically finds lung nodules in CT scans using 3D deep learning helping doctors detect lung cancer earlier and more accurately.",
   size=14, color=TEXT2, align=PP_ALIGN.CENTER)

tb(s, Inches(4.5), Inches(3.5), Inches(4), Inches(0.3), "PROJECT GUIDE", size=11, bold=True, color=TEXT3, align=PP_ALIGN.CENTER)
add_card(s, Inches(5), Inches(3.85), Inches(3), Inches(0.8), "J. Ravindra Babu", "Assistant Professor", title_size=14, body_size=11)

tb(s, Inches(1), Inches(4.9), Inches(11), Inches(0.3), "TEAM MEMBERS", size=11, bold=True, color=TEXT3, align=PP_ALIGN.CENTER)
members = ["B. Sandeep Raghavendra", "A. Jaswanth Kumar", "G. Vignan", "J. Ganesh", "K. Madhu"]
start_x = 1.3
for i, name in enumerate(members):
    add_card(s, Inches(start_x + i * 2.2), Inches(5.3), Inches(2.0), Inches(0.7), name, "", title_size=11, body_size=10)

# ============================================================
# SLIDE 2 — ABSTRACT
# ============================================================
s = prs.slides.add_slide(BLANK); add_bg(s)
add_accent_header(s, 2, "Abstract")

abstract = (
    "Lung cancer is the leading cause of cancer deaths worldwide. Early detection of small growths "
    "called nodules in CT scans can significantly improve survival rates. However, manually reviewing "
    "hundreds of scan slices is slow and prone to human error.\n\n"
    "This project presents a complete web-based application that uses a 3D deep learning model to "
    "automatically detect and locate lung nodules in CT scans. The system processes full 3D scan volumes, "
    "identifies suspicious regions, and generates diagnostic reports — acting as an AI assistant for radiologists.\n\n"
    "Trained on the LUNA16 benchmark dataset (888 CT scans), our model achieves 94.2% detection rate "
    "with only 1.79 false alarms per scan, making it suitable as a clinical screening aid."
)
shape = slide_shapes = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(1.5), Inches(10), Inches(5.0))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0xF8, 0xFA, 0xFC)
shape.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
shape.line.width = Pt(1)
tf = shape.text_frame
tf.word_wrap = True
tf.margin_left = Pt(20)
tf.margin_right = Pt(20)
tf.margin_top = Pt(20)
p = tf.paragraphs[0]
p.text = abstract
p.font.size = Pt(16)
p.font.color.rgb = TEXT2
p.font.name = FONT
p.line_spacing = Pt(26)

# ============================================================
# SLIDE 3 — OBJECTIVES
# ============================================================
s = prs.slides.add_slide(BLANK); add_bg(s)
add_accent_header(s, 3, "Objectives")

objectives = [
    ("🎯 Automated Detection", "Build a system that can automatically find lung nodules in 3D CT scans without manual intervention."),
    ("📊 High Accuracy", "Achieve over 90% detection rate while keeping false alarms low — reducing missed diagnoses."),
    ("🌐 Web Application", "Create an easy-to-use web interface where doctors can upload scans and view results instantly."),
    ("📋 Report Generation", "Automatically generate diagnostic reports with nodule locations, sizes, and confidence scores."),
    ("👁️ Visual Inspection", "Provide an interactive slice viewer so doctors can scroll through scans and verify AI findings."),
    ("🏥 Clinical Aid", 'Serve as a "second reader" that assists radiologists — not replaces them — in faster screening.'),
]
positions = [
    (1.0, 1.5), (6.8, 1.5),
    (1.0, 3.1), (6.8, 3.1),
    (1.0, 4.7), (6.8, 4.7),
]
for (title, body), (x, y) in zip(objectives, positions):
    add_card(s, Inches(x), Inches(y), Inches(5.5), Inches(1.3), title, body)

# ============================================================
# SLIDE 4 — SYSTEM OVERVIEW
# ============================================================
s = prs.slides.add_slide(BLANK); add_bg(s)
add_accent_header(s, 4, "System Overview")
tb(s, Inches(2), Inches(1.3), Inches(9), Inches(0.5),
   "The application has three main parts working together: a user interface, a server, and an AI engine.",
   size=14, color=TEXT2, align=PP_ALIGN.CENTER)

# Architecture boxes
boxes = [("🖥️ Frontend", "React Web Interface"), ("⚙️ Backend Server", "Flask REST API"), ("🧠 AI Engine", "3D U-Net Model")]
colors = [SKY, ACCENT, TEAL]
for i, ((title, sub), clr) in enumerate(zip(boxes, colors)):
    x = Inches(1.5 + i * 4.0)
    shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.0), Inches(3.2), Inches(1.0))
    shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor(0xF8, 0xFA, 0xFC)
    shape.line.color.rgb = clr; shape.line.width = Pt(2)
    tf = shape.text_frame; tf.word_wrap = True
    tf.margin_left = Pt(10); tf.margin_top = Pt(8)
    p = tf.paragraphs[0]; p.text = title; p.font.size = Pt(15); p.font.bold = True; p.font.name = FONT; p.font.color.rgb = TEXT; p.alignment = PP_ALIGN.CENTER
    add_para(tf, sub, size=12, color=TEXT3, align=PP_ALIGN.CENTER)

# Arrows
for i in range(2):
    tb(s, Inches(4.7 + i * 4.0), Inches(2.25), Inches(0.5), Inches(0.5), "⇄", size=24, color=ACCENT, align=PP_ALIGN.CENTER)

# Two cards
left_items = ["Home — Overview and quick access", "Analyze — Upload scans and view AI results", "Performance — Model accuracy charts", "About — Team and project information"]
right_items = ["Upload handler accepts CT scan files", "Preprocessing normalizes the 3D volume", "AI model scans the entire volume", "Results saved to database and sent to user"]

card_l = add_card(s, Inches(1.0), Inches(3.4), Inches(5.5), Inches(3.5), "What the User Sees", "", body_size=13)
tf = card_l.text_frame
for item in left_items:
    add_bullet(tf, f"▸ {item}", size=13, color=TEXT2)

card_r = add_card(s, Inches(6.8), Inches(3.4), Inches(5.5), Inches(3.5), "What Happens Behind", "", body_size=13)
tf = card_r.text_frame
for item in right_items:
    add_bullet(tf, f"▸ {item}", size=13, color=TEXT2)

# ============================================================
# SLIDE 5 — WORKFLOW
# ============================================================
s = prs.slides.add_slide(BLANK); add_bg(s)
add_accent_header(s, 5, "Workflow")
tb(s, Inches(2), Inches(1.3), Inches(9), Inches(0.5),
   "Step-by-step flow of how a CT scan goes from upload to diagnosis report.",
   size=14, color=TEXT2, align=PP_ALIGN.CENTER)

row1 = [("1. Upload", "CT scan file"), ("2. Load Scan", "Read 3D volume"), ("3. Preprocess", "Normalize & prepare"),
        ("4. Split Patches", "Divide into blocks"), ("5. AI Analysis", "3D U-Net prediction")]
row2 = [("6. Combine", "Merge predictions"), ("7. Find Nodules", "Locate suspects"), ("8. Score", "Assign confidence"),
        ("9. Save Report", "Store in database"), ("10. Display", "Show to doctor")]

for row_idx, row in enumerate([row1, row2]):
    y = Inches(2.0 + row_idx * 2.2)
    for i, (title, sub) in enumerate(row):
        x = Inches(0.5 + i * 2.5)
        shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.0), Inches(1.0))
        shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor(0xF8, 0xFA, 0xFC)
        shape.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0); shape.line.width = Pt(1)
        tf = shape.text_frame; tf.word_wrap = True; tf.margin_left = Pt(8); tf.margin_top = Pt(8)
        p = tf.paragraphs[0]; p.text = title; p.font.size = Pt(13); p.font.bold = True; p.font.name = FONT; p.font.color.rgb = TEXT; p.alignment = PP_ALIGN.CENTER
        add_para(tf, sub, size=11, color=TEXT3, align=PP_ALIGN.CENTER)
        if i < 4:
            tb(s, Inches(2.5 + i * 2.5), y + Inches(0.25), Inches(0.5), Inches(0.5), "→", size=20, color=ACCENT, align=PP_ALIGN.CENTER)

add_card(s, Inches(3.5), Inches(6.0), Inches(6), Inches(0.7), "Total processing time: ~12 seconds per full scan volume", "", title_size=14)

# ============================================================
# SLIDE 6 — TECHNOLOGIES USED
# ============================================================
s = prs.slides.add_slide(BLANK); add_bg(s)
add_accent_header(s, 6, "Technologies Used")

techs = [
    ("🧠 PyTorch", "Deep learning framework used to build and train the 3D neural network model."),
    ("⚛️ React", "JavaScript library for building the interactive web user interface."),
    ("🌐 Flask", "Python web framework that serves as the backend server handling requests."),
    ("⚡ Vite", "Modern build tool for fast development and optimized production builds."),
    ("🗄️ SQLite", "Lightweight database for storing patient records and diagnostic reports."),
    ("📐 TypeScript", "Typed JavaScript for more reliable and maintainable frontend code."),
]
positions = [
    (0.8, 1.5), (4.6, 1.5), (8.4, 1.5),
    (0.8, 3.7), (4.6, 3.7), (8.4, 3.7),
]
for (title, body), (x, y) in zip(techs, positions):
    add_card(s, Inches(x), Inches(y), Inches(3.5), Inches(1.8), title, body)

# ============================================================
# SLIDE 7 — DATASETS & TOOLS
# ============================================================
s = prs.slides.add_slide(BLANK); add_bg(s)
add_accent_header(s, 7, "Datasets & Tools")

# LUNA16 card
shape = add_card(s, Inches(0.8), Inches(1.5), Inches(5.5), Inches(5.3), "LUNA16 Dataset", "")
tf = shape.text_frame
add_para(tf, "The standard benchmark dataset for lung nodule detection research, used worldwide.", size=14)
add_para(tf, "")
add_para(tf, "888 CT Scans  •  1,186 Nodule Annotations", size=16, bold=True, color=TEAL)
add_para(tf, "")
add_para(tf, "Data Format:", size=14, bold=True, color=TEXT)
add_bullet(tf, "▸ .mhd — Metadata file (scan info)", size=13, color=TEXT2)
add_bullet(tf, "▸ .raw — Actual 3D image data", size=13, color=TEXT2)
add_bullet(tf, "▸ Each scan has 200–400 slices", size=13, color=TEXT2)

# Tools
tools = [
    ("🏥 SimpleITK", "Medical imaging library for reading and processing CT scan files."),
    ("📊 NumPy & SciPy", "Scientific computing libraries for handling large 3D arrays."),
    ("📈 Recharts", "Charting library for performance graphs and metrics in the frontend."),
]
for i, (title, body) in enumerate(tools):
    add_card(s, Inches(6.8), Inches(1.5 + i * 1.9), Inches(5.5), Inches(1.5), title, body)

# ============================================================
# SLIDE 8 — PREPROCESSING STRATEGIES
# ============================================================
s = prs.slides.add_slide(BLANK); add_bg(s)
add_accent_header(s, 8, "Preprocessing Strategies")
tb(s, Inches(2), Inches(1.3), Inches(9), Inches(0.5),
   "Before feeding scans to the AI model, the raw data needs to be cleaned and prepared.",
   size=14, color=TEXT2, align=PP_ALIGN.CENTER)

preps = [
    ("🔢 Intensity Normalization", "CT scan values are converted to a 0–1 scale. Air and bone are clipped to focus on lung tissue."),
    ("✂️ Patch Extraction", "Full scans are cut into small 3D blocks (32×64×64 pixels) the model can process one at a time."),
    ("🔄 Data Augmentation", "Training patches are randomly flipped, rotated, and have noise added for better generalization."),
    ("⚖️ Class Balancing", 'Equal numbers of "nodule" and "non-nodule" patches prevent model bias.'),
]
positions = [(0.8, 2.0), (6.8, 2.0), (0.8, 4.0), (6.8, 4.0)]
for (title, body), (x, y) in zip(preps, positions):
    add_card(s, Inches(x), Inches(y), Inches(5.5), Inches(1.7), title, body)

# Flow
flow = ["Raw CT Scan", "Normalize", "Extract Patches", "Augment", "Ready for AI"]
for i, label in enumerate(flow):
    x = Inches(0.5 + i * 2.5)
    shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(6.0), Inches(2.0), Inches(0.7))
    shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor(0xF8, 0xFA, 0xFC)
    shape.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0); shape.line.width = Pt(1)
    tf = shape.text_frame; tf.word_wrap = True; tf.margin_top = Pt(6)
    p = tf.paragraphs[0]; p.text = label; p.font.size = Pt(12); p.font.bold = True; p.font.name = FONT; p.font.color.rgb = TEXT; p.alignment = PP_ALIGN.CENTER
    if i < 4:
        tb(s, Inches(2.5 + i * 2.5), Inches(6.05), Inches(0.5), Inches(0.5), "→", size=18, color=ACCENT, align=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 9 — ARCHITECTURE
# ============================================================
s = prs.slides.add_slide(BLANK); add_bg(s)
add_accent_header(s, 9, "Architecture of the Project")
tb(s, Inches(1.5), Inches(1.3), Inches(10), Inches(0.5),
   "The AI model uses a 3D U-Net design — an encoder-decoder network that learns to highlight nodules.",
   size=14, color=TEXT2, align=PP_ALIGN.CENTER)

# Left column
left = [
    ("⬇️ Encoder (Compress)", "The scan passes through 4 stages that progressively shrink the image while extracting features.", ACCENT),
    ("🔄 Bottleneck (Core)", "The deepest layer captures the most meaningful patterns about lung tissue vs. nodules.", AMBER),
    ("⬆️ Decoder (Expand)", "Expands back to original size using skip connections to recover fine details.", TEAL),
]
for i, (title, body, clr) in enumerate(left):
    c = add_card(s, Inches(0.8), Inches(2.0 + i * 1.7), Inches(5.5), Inches(1.4), title, body)
    c.line.color.rgb = clr

# Right column
right = [
    ("🧱 Residual Blocks", "Shortcut connections help the model train deeper without losing information."),
    ("🔗 Skip Connections", "Links between encoder and decoder preserve fine spatial details."),
    ("⚡ Mixed Precision", "Uses 16-bit and 32-bit math to train 2× faster with half the memory."),
]
for i, (title, body) in enumerate(right):
    add_card(s, Inches(6.8), Inches(2.0 + i * 1.7), Inches(5.5), Inches(1.4), title, body)

# ============================================================
# SLIDE 10 — ADVANTAGES
# ============================================================
s = prs.slides.add_slide(BLANK); add_bg(s)
add_accent_header(s, 10, "Advantages")

advs = [
    ("🔍 3D Analysis", "Analyzes full 3D volume — capturing spatial relationships between slices that flat images miss."),
    ("⚡ Fast Results", "Processes a complete CT scan in ~12 seconds vs. 15–20 minutes for manual review."),
    ("📊 High Detection Rate", "94.2% sensitivity with only 1.79 false alarms per scan."),
    ("🤝 Doctor-Friendly", "Designed as a helpful assistant. Doctors verify every AI finding through the slice viewer."),
    ("💻 Runs Anywhere", "Supports NVIDIA GPUs, Apple Silicon, and regular CPUs — no specialized hardware required."),
    ("🌐 Web-Based", "No software installation needed. Works through any modern web browser."),
]
positions = [(0.8, 1.5), (6.8, 1.5), (0.8, 3.3), (6.8, 3.3), (0.8, 5.1), (6.8, 5.1)]
for (title, body), (x, y) in zip(advs, positions):
    add_card(s, Inches(x), Inches(y), Inches(5.5), Inches(1.5), title, body)

# ============================================================
# SLIDE 11 — APPLICATIONS
# ============================================================
s = prs.slides.add_slide(BLANK); add_bg(s)
add_accent_header(s, 11, "Applications")

apps = [
    ("🏥 Hospital Screening", "First-pass screening tool to quickly clear normal scans and flag suspicious ones."),
    ("🩺 Lung Cancer Screening", "Supports large-scale population screening with rapid, consistent assessment."),
    ("🏫 Medical Education", "Teaching tool for medical students to understand nodule characteristics and AI diagnosis."),
    ("🔬 Research Platform", "Foundation for researchers to test new AI models on standard benchmarks."),
    ("🌍 Rural Healthcare", "Brings expert-level screening to areas with limited access to specialist radiologists."),
    ("⏰ Emergency Triage", "Quickly identifies scans needing immediate attention in busy emergency departments."),
]
positions = [(0.8, 1.5), (6.8, 1.5), (0.8, 3.3), (6.8, 3.3), (0.8, 5.1), (6.8, 5.1)]
for (title, body), (x, y) in zip(apps, positions):
    add_card(s, Inches(x), Inches(y), Inches(5.5), Inches(1.5), title, body)

# ============================================================
# SLIDE 12 — LIMITATIONS
# ============================================================
s = prs.slides.add_slide(BLANK); add_bg(s)
add_accent_header(s, 12, "Limitations")

lims = [
    ("📏 Small Nodule Challenge", "Nodules < 3mm may be missed as the model is trained on annotations ≥ 3mm.", AMBER),
    ("💾 Memory Requirements", "3D data processing requires significant RAM. Training needs a powerful GPU.", ROSE),
    ("📂 File Format Dependency", "Currently only supports MetaImage (.mhd/.raw). DICOM not yet supported.", SKY),
    ("🔬 No Malignancy Classification", "Detects nodules but does not classify them as benign or malignant.", ACCENT),
    ("🌐 Single Dataset Training", "Trained only on LUNA16. Performance may vary on different scanners.", TEAL),
    ("⚕️ Not Standalone Diagnostic", "Requires expert radiologist verification for clinical decision-making.", EMERALD),
]
positions = [(0.8, 1.5), (6.8, 1.5), (0.8, 3.3), (6.8, 3.3), (0.8, 5.1), (6.8, 5.1)]
for (title, body, clr), (x, y) in zip(lims, positions):
    c = add_card(s, Inches(x), Inches(y), Inches(5.5), Inches(1.5), title, body)
    c.line.color.rgb = clr

# ============================================================
# SLIDE 13 — UML DIAGRAMS
# ============================================================
s = prs.slides.add_slide(BLANK); add_bg(s)
add_accent_header(s, 13, "UML Diagrams")

uml1 = os.path.join(DIR, "UML-1.png")
uml2 = os.path.join(DIR, "UML-2.png")

if os.path.exists(uml1):
    tb(s, Inches(0.8), Inches(1.4), Inches(5.5), Inches(0.4), "System Workflow Diagram", size=16, bold=True, color=TEXT, align=PP_ALIGN.CENTER)
    s.shapes.add_picture(uml1, Inches(0.8), Inches(1.9), Inches(5.5), Inches(4.5))
else:
    tb(s, Inches(1), Inches(2), Inches(5), Inches(1), "[UML-1.png not found]", size=14, color=TEXT3)

if os.path.exists(uml2):
    tb(s, Inches(6.8), Inches(1.4), Inches(5.5), Inches(0.4), "Performance Comparison Diagram", size=16, bold=True, color=TEXT, align=PP_ALIGN.CENTER)
    s.shapes.add_picture(uml2, Inches(6.8), Inches(1.9), Inches(5.5), Inches(4.5))
else:
    tb(s, Inches(7), Inches(2), Inches(5), Inches(1), "[UML-2.png not found]", size=14, color=TEXT3)

# ============================================================
# SLIDE 14 — CONCLUSION
# ============================================================
s = prs.slides.add_slide(BLANK); add_bg(s)
add_accent_header(s, 14, "Conclusion")

conclusion = (
    "We have successfully developed an end-to-end deep learning system for automated lung nodule detection. "
    "The system combines a powerful 3D AI model with an intuitive web interface, making advanced screening "
    "technology accessible and easy to use."
)
shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(1.5), Inches(11), Inches(1.5))
shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor(0xF8, 0xFA, 0xFC)
shape.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0); shape.line.width = Pt(1)
tf = shape.text_frame; tf.word_wrap = True; tf.margin_left = Pt(16); tf.margin_top = Pt(12)
p = tf.paragraphs[0]; p.text = conclusion; p.font.size = Pt(15); p.font.color.rgb = TEXT2; p.font.name = FONT; p.line_spacing = Pt(24)

# Metrics
metrics = [("94.2%", "Detection Rate", TEAL), ("1.79", "False Positives / Scan", ACCENT), ("~12s", "Processing Time", AMBER)]
for i, (val, label, clr) in enumerate(metrics):
    x = Inches(1.0 + i * 3.8)
    shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(3.3), Inches(3.3), Inches(1.2))
    shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor(0xF8, 0xFA, 0xFC)
    shape.line.color.rgb = clr; shape.line.width = Pt(2)
    tf = shape.text_frame; tf.word_wrap = True; tf.margin_top = Pt(10)
    p = tf.paragraphs[0]; p.text = val; p.font.size = Pt(28); p.font.bold = True; p.font.color.rgb = clr; p.font.name = FONT; p.alignment = PP_ALIGN.CENTER
    add_para(tf, label, size=12, color=TEXT3, align=PP_ALIGN.CENTER)

tb(s, Inches(1.0), Inches(4.8), Inches(5), Inches(0.4), "Future Scope", size=18, bold=True, color=AMBER)
future = [
    ("🔮 Malignancy Classification", "Benign vs malignant nodule prediction"),
    ("☁️ Cloud Deployment", "DICOM format support for hospital integration"),
    ("🧠 Explainable AI", "Heatmaps showing why the model flagged a region"),
    ("🔒 Privacy-Preserving", "Training across multiple hospitals"),
]
positions = [(1.0, 5.3), (6.8, 5.3), (1.0, 6.2), (6.8, 6.2)]
for (title, body), (x, y) in zip(future, positions):
    add_card(s, Inches(x), Inches(y), Inches(5.5), Inches(0.8), title, body, title_size=13, body_size=11)

# ============================================================
# SLIDE 15 — THANK YOU
# ============================================================
s = prs.slides.add_slide(BLANK); add_bg(s)
tb(s, Inches(1), Inches(1.5), Inches(11), Inches(1.0), "Thank You!", size=44, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
tb(s, Inches(1), Inches(2.6), Inches(11), Inches(0.5), "Questions & Discussion", size=20, color=TEXT2, align=PP_ALIGN.CENTER)

shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.5), Inches(3.4), Inches(6), Inches(1.6))
shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor(0xF8, 0xFA, 0xFC)
shape.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0); shape.line.width = Pt(1)
tf = shape.text_frame; tf.word_wrap = True; tf.margin_left = Pt(16); tf.margin_top = Pt(14)
p = tf.paragraphs[0]; p.text = "Volumetric CT Scan Analysis for Pulmonary Nodule Detection"; p.font.size = Pt(16); p.font.bold = True; p.font.name = FONT; p.font.color.rgb = TEXT; p.alignment = PP_ALIGN.CENTER
add_para(tf, "B.Tech Final Year Project — CSE (Data Science)", size=13, color=TEXT2, align=PP_ALIGN.CENTER)
add_para(tf, "KKR & KSR Institute of Technology and Sciences, Guntur", size=12, color=TEXT3, align=PP_ALIGN.CENTER)

# Member tags
for i, name in enumerate(members):
    x = Inches(1.6 + i * 2.1)
    shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(5.5), Inches(1.9), Inches(0.5))
    shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor(0xEE, 0xF2, 0xFF)
    shape.line.fill.background()
    tf = shape.text_frame; tf.margin_top = Pt(4)
    p = tf.paragraphs[0]; p.text = name; p.font.size = Pt(11); p.font.bold = True; p.font.color.rgb = ACCENT; p.font.name = FONT; p.alignment = PP_ALIGN.CENTER

tb(s, Inches(1), Inches(6.2), Inches(11), Inches(0.4),
   "Guide: J. Ravindra Babu, Assistant Professor", size=13, color=TEXT3, align=PP_ALIGN.CENTER)

# ---------- save ----------
out = os.path.join(DIR, "presentation.pptx")
prs.save(out)
print(f"✅ Saved to {out}")
